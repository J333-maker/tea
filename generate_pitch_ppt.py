from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def main() -> None:
    base = Path(__file__).resolve().parent
    out = base / "三创赛路演答辩PPT模板_安顶山云雾茶数字焕新平台.pptx"

    col_primary = RGBColor(0x2B, 0x7A, 0x3B)
    col_secondary = RGBColor(0x4C, 0xAF, 0x50)
    col_dark = RGBColor(0x1B, 0x4D, 0x24)
    col_light = RGBColor(0xE8, 0xF5, 0xE9)
    col_text = RGBColor(0x22, 0x22, 0x22)

    img_logo = base / "logo.jpg"
    img_cover = base / "封面.jpg"
    img_mascot = base / "茶小泽2.png"
    img_overview = base / "总览.jpg"
    img_guide = base / "导览.jpg"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_top_bar(slide, title: str, subtitle: str | None = None, show_logo: bool = True) -> None:
        w = prs.slide_width
        bar_h = Inches(0.85)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w, bar_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = col_primary
        bar.line.fill.background()

        tx = slide.shapes.add_textbox(Inches(0.65), Inches(0.17), w - Inches(2.2), Inches(0.6))
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.name = "Microsoft YaHei"
            p2.font.size = Pt(12)
            p2.font.color.rgb = RGBColor(0xF1, 0xF8, 0xF4)

        if show_logo and img_logo.exists():
            slide.shapes.add_picture(str(img_logo), w - Inches(1.55), Inches(0.16), height=Inches(0.55))

    def add_section_card(slide, x, y, w, h, title: str, bullets: list[str]) -> None:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        card.line.width = Pt(1)

        head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.42))
        head.fill.solid()
        head.fill.fore_color.rgb = col_light
        head.line.fill.background()

        tb = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.06), w - Inches(0.44), Inches(0.32))
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col_dark

        bt = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.55), w - Inches(0.44), h - Inches(0.65))
        tf2 = bt.text_frame
        tf2.clear()
        for i, b in enumerate(bullets):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p.text = b
            p.level = 0
            p.font.name = "Microsoft YaHei"
            p.font.size = Pt(14)
            p.font.color.rgb = col_text
            p.space_after = Pt(6)

    def add_kpi_row(slide, y, items: list[tuple[str, str]]) -> None:
        w = prs.slide_width
        gap = Inches(0.25)
        left = Inches(0.65)
        usable = w - left * 2
        box_w = (usable - gap * (len(items) - 1)) / len(items)
        box_h = Inches(1.2)
        for i, (k, v) in enumerate(items):
            x = left + i * (box_w + gap)
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            box.line.color.rgb = RGBColor(0xD9, 0xE8, 0xDD)
            box.line.width = Pt(1)
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.12), box_h)
            accent.fill.solid()
            accent.fill.fore_color.rgb = col_secondary
            accent.line.fill.background()

            t1 = slide.shapes.add_textbox(x + Inches(0.24), y + Inches(0.15), box_w - Inches(0.34), Inches(0.35))
            p1 = t1.text_frame.paragraphs[0]
            p1.text = k
            p1.font.name = "Microsoft YaHei"
            p1.font.size = Pt(12)
            p1.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

            t2 = slide.shapes.add_textbox(x + Inches(0.24), y + Inches(0.5), box_w - Inches(0.34), Inches(0.6))
            p2 = t2.text_frame.paragraphs[0]
            p2.text = v
            p2.font.name = "Microsoft YaHei"
            p2.font.size = Pt(24)
            p2.font.bold = True
            p2.font.color.rgb = col_primary

    slides = []

    slide = prs.slides.add_slide(blank)
    if img_cover.exists():
        slide.shapes.add_picture(str(img_cover), 0, 0, width=prs.slide_width, height=prs.slide_height)
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0x00, 0x30, 0x14)
        overlay.fill.transparency = 0.55
        overlay.line.fill.background()
    else:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = col_primary
        bg.line.fill.background()

    if img_logo.exists():
        slide.shapes.add_picture(str(img_logo), Inches(0.9), Inches(0.7), height=Inches(0.7))
    if img_mascot.exists():
        slide.shapes.add_picture(str(img_mascot), prs.slide_width - Inches(3.2), Inches(3.2), width=Inches(2.4))

    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(10.8), Inches(2.2))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "安顶山云雾茶数字焕新平台"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p2 = tf.add_paragraph()
    p2.text = "三创赛路演答辩 PPT 模板（绿色主题｜20页）"
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xE6, 0xFF, 0xEF)
    p2.space_before = Pt(14)
    p3 = tf.add_paragraph()
    p3.text = "团队：__________    学校：__________    日期：____/__/__"
    p3.font.name = "Microsoft YaHei"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0xF3, 0xFF, 0xF7)
    p3.space_before = Pt(18)
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_top_bar(slide, "目录", "路演逻辑：痛点 → 方案 → 产品 → 商业 → 运营 → 团队")
    add_section_card(slide, Inches(0.75), Inches(1.35), Inches(5.9), Inches(5.6), "Part 1｜项目概述", ["赛道与背景", "目标用户与核心痛点", "解决方案定位"])
    add_section_card(slide, Inches(6.85), Inches(1.35), Inches(5.9), Inches(5.6), "Part 2｜产品与技术", ["平台功能地图（B端+ C端）", "AI 客服（茶小泽）与知识问答", "数据闭环与可视化"])
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_top_bar(slide, "项目背景与机会")
    add_kpi_row(slide, Inches(1.35), [("产业现状", "分散"), ("用户诉求", "透明"), ("经营目标", "提效"), ("品牌方向", "数字化")])
    add_section_card(
        slide,
        Inches(0.75),
        Inches(2.75),
        Inches(12.58),
        Inches(4.1),
        "背景要点（可替换为调研数据）",
        [
            "茶叶经营链条长：种植/加工/流通/零售/文旅，信息割裂导致协同成本高",
            "消费者关注溯源、等级、口感与性价比，但线上内容碎片化、决策困难",
            "景区与文旅业态需要“导览+预约+消费”一体化入口，提升转化与复购",
            "数据资产沉淀不足：订单、客户、产品与活动数据难以实时分析与决策",
        ],
    )
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_top_bar(slide, "核心痛点拆解")
    add_section_card(
        slide,
        Inches(0.75),
        Inches(1.35),
        Inches(6.1),
        Inches(5.5),
        "消费者侧（C端）",
        ["信息不对称：等级/产地/工艺/价格难以快速理解与比较", "决策链长：导览、预约、购买、售后入口分散", "服务压力：咨询高频且重复，人工响应不稳定"],
    )
    add_section_card(
        slide,
        Inches(6.95),
        Inches(1.35),
        Inches(5.9),
        Inches(5.5),
        "经营侧（B端）",
        ["数据滞后：销量、订单、客单价、热销品无法实时看", "运营缺抓手：活动效果难评估，复购与会员运营难闭环", "渠道协同弱：线上线下活动与库存、履约信息难统一"],
    )
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_top_bar(slide, "解决方案：一体化平台 + 数据闭环")
    add_section_card(
        slide,
        Inches(0.75),
        Inches(1.35),
        Inches(12.58),
        Inches(2.35),
        "一句话定位",
        ["打造“导览预约 + 交易中心 + 产业大脑 + AI 客服”的茶产业互联网综合服务平台，连接消费者与产业运营。"],
    )
    add_section_card(
        slide,
        Inches(0.75),
        Inches(3.9),
        Inches(12.58),
        Inches(2.95),
        "价值主张",
        ["对用户：更快获得可信答案（溯源/等级/路线/预约/购买），更顺畅完成转化", "对商家：用数据驱动选品、定价、活动与供应链，降低客服与运营成本", "对品牌：统一内容与服务口径，构建数字资产与长期复购体系"],
    )
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_top_bar(slide, "产品总览（对应网站模块）", "C端门户 + B端看板 + 数据/AI能力")
    if img_overview.exists():
        slide.shapes.add_picture(str(img_overview), Inches(0.75), Inches(1.45), width=Inches(6.1), height=Inches(5.7))
    add_section_card(
        slide,
        Inches(6.95),
        Inches(1.45),
        Inches(5.88),
        Inches(5.7),
        "功能地图",
        [
            "首页：平台定位与核心入口（导览/预约/交易）",
            "茶文旅：项目预约、导览图、活动推荐",
            "茶产业：溯源与产业链协同、数字资产沉淀",
            "茶科技：数据分析、品质与工艺展示",
            "茶服务：团购定制、会员与售后",
            "茶碳汇：绿色价值与生态叙事（可扩展）",
            "交易中心：茶叶/文创交易与履约流程",
            "产业大脑：经营数据看板与决策辅助",
            "AI 客服：茶小泽，支持FAQ与大模型增强问答",
        ],
    )
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_top_bar(slide, "用户旅程与业务闭环")
    add_section_card(
        slide,
        Inches(0.75),
        Inches(1.35),
        Inches(12.58),
        Inches(5.5),
        "从“看见”到“复购”的闭环（示例）",
        ["触达：内容/活动/导览入口 → 进入平台", "咨询：茶小泽解答（导览/预约/溯源/就餐/团购）", "转化：交易中心下单/预约项目 → 支付/履约", "沉淀：订单、客户、产品、行为数据进入数据库", "运营：看板分析 → 活动投放/选品/价格策略 → 复购与口碑"],
    )
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_top_bar(slide, "AI 客服「茶小泽」：从FAQ到大模型增强")
    if img_mascot.exists():
        slide.shapes.add_picture(str(img_mascot), Inches(0.9), Inches(1.65), width=Inches(3.0))
    add_section_card(
        slide,
        Inches(4.2),
        Inches(1.35),
        Inches(8.98),
        Inches(2.75),
        "能力说明",
        ["本地FAQ：覆盖高频问题，秒级响应，离线可用", "大模型增强：接入兼容 OpenAI 的大模型 API，支持更复杂的追问与多轮对话", "安全策略：不在前端暴露密钥，服务端转发调用，失败自动回退 FAQ"],
    )
    add_section_card(
        slide,
        Inches(4.2),
        Inches(4.25),
        Inches(8.98),
        Inches(2.6),
        "可展示Demo点",
        ["“导览路线怎么走？”→ 给路线 + 引导查看导览图", "“特选级和精选级怎么选？”→ 解释等级差异 + 推荐场景", "“团购定制怎么做？”→ 给流程模板（建议替换为真实政策/数据）"],
    )
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_top_bar(slide, "技术架构（可按实际实现调整）")
    add_section_card(
        slide,
        Inches(0.75),
        Inches(1.35),
        Inches(12.58),
        Inches(5.5),
        "架构示意",
        [
            "前端：HTML5 + Bootstrap + 图表（Chart.js）+ 交互组件",
            "后端：Python 轻量服务（静态资源 + API），SQLite 存储业务数据",
            "数据：脚本导出 data.json，支撑看板与运营分析",
            "AI：/api/chat 服务端转发到大模型（环境变量配置 Key/模型/网关）",
            "部署：本地/云服务器均可，后续可升级为 Flask/FastAPI/容器化",
        ],
    )
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_top_bar(slide, "产业大脑（B端）：指标体系与决策支持")
    add_kpi_row(slide, Inches(1.35), [("GMV", "____"), ("订单数", "____"), ("客单价", "____"), ("活跃客户", "____")])
    add_section_card(
        slide,
        Inches(0.75),
        Inches(2.75),
        Inches(12.58),
        Inches(4.1),
        "看板模块（示例）",
        ["关键指标：销售额、订单、客单价、客户活跃", "趋势分析：月度销售趋势、品类占比、热销排行", "运营辅助：最新订单、异常订单/退换货（可扩展）", "决策动作：选品/补货/活动/渠道投放建议（结合AI可增强）"],
    )
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_top_bar(slide, "交易中心（C端）：茶品 + 文创 + 预约")
    add_section_card(slide, Inches(0.75), Inches(1.35), Inches(6.1), Inches(5.5), "核心交易", ["茶叶预定：按等级/规格下单，支持礼盒与团购", "文创周边：吉祥物系列、定制茶具、联名产品", "文旅预约：采茶/炒茶/研学，形成“体验→购买”转化"])
    if img_guide.exists():
        slide.shapes.add_picture(str(img_guide), Inches(7.05), Inches(1.55), width=Inches(5.98), height=Inches(5.1))
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_top_bar(slide, "竞品分析与差异化")
    add_section_card(
        slide,
        Inches(0.75),
        Inches(1.35),
        Inches(12.58),
        Inches(5.5),
        "差异化（示例模板，可替换为你的真实竞品）",
        ["平台一体化：导览/预约/交易/数据看板/AI客服在同一入口，减少跳转流失", "数据闭环：订单与运营数据沉淀，支持选品、活动复盘与增长", "内容可信：溯源与工艺讲解标准化，降低理解门槛", "服务可规模化：AI客服承担高频咨询，人工专注高价值问题"],
    )
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_top_bar(slide, "商业模式（收益来源）")
    add_section_card(slide, Inches(0.75), Inches(1.35), Inches(6.1), Inches(5.5), "收入", ["交易佣金/毛利：茶叶、文创、团购礼盒", "服务费：企业团购定制、活动策划、供应链协同", "SaaS/订阅（可选）：给合作茶企提供看板与运营工具", "品牌合作：联名、文旅套餐、渠道分销"])
    add_section_card(slide, Inches(6.95), Inches(1.35), Inches(5.9), Inches(5.5), "成本", ["研发与运维：服务器、域名、维护与迭代", "内容与运营：素材、活动、客服与售后", "履约与供应链：仓配、品控与退换货", "营销获客：渠道投放、合作分成"])
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_top_bar(slide, "目标用户与使用场景")
    add_section_card(slide, Inches(0.75), Inches(1.35), Inches(12.58), Inches(5.5), "用户画像（模板）", ["游客/体验者：想要路线、预约、就餐、纪念品推荐", "茶叶爱好者：关注等级差异、口感风味、性价比与溯源", "企业团购：关注定制、交付、预算与品牌表达", "茶企/运营者：关注销售、库存、活动效果、复购与客群结构"])
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_top_bar(slide, "运营策略：内容 + 活动 + 私域")
    add_section_card(slide, Inches(0.75), Inches(1.35), Inches(12.58), Inches(5.5), "策略拆解（可替换）", ["内容：等级科普/冲泡指南/工艺故事/溯源展示，统一口径", "活动：节气上新、采茶季、研学营、企业团购季，形成节点转化", "私域：会员积分、复购券、团购名单、售后关怀", "增长：短视频/小红书/抖音/线下景区导流 → 平台承接"])
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_top_bar(slide, "落地计划（里程碑）")
    add_section_card(slide, Inches(0.75), Inches(1.35), Inches(12.58), Inches(5.5), "三阶段推进（模板）", ["阶段1：内容与交易上线（茶品/文创/预约），跑通支付与履约", "阶段2：数据看板与运营体系（指标、复盘、活动模板）", "阶段3：AI增强与智能化（多轮客服、推荐、运营策略生成）"])
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_top_bar(slide, "数据治理与安全合规")
    add_section_card(slide, Inches(0.75), Inches(1.35), Inches(12.58), Inches(5.5), "要点（模板）", ["权限：B端账号分级，关键操作留痕（可扩展）", "隐私：不在前端暴露API密钥；用户数据最小化采集", "稳定：接口失败自动降级（AI→FAQ），保障核心流程可用", "审计：关键指标可追溯，便于复盘与监管检查"])
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_top_bar(slide, "财务预测（模板）")
    add_section_card(slide, Inches(0.75), Inches(1.35), Inches(12.58), Inches(5.5), "表格建议（填空）", ["收入：茶品销售____、文创销售____、团购定制____、服务/订阅____", "成本：履约____、营销____、研发运维____、人工____", "关键指标：CAC、复购率、毛利率、ARPU、库存周转（建议在答辩中重点解释）"])
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_top_bar(slide, "团队介绍与分工")
    add_section_card(slide, Inches(0.75), Inches(1.35), Inches(6.1), Inches(5.5), "团队成员（填空）", ["队长：__________（产品/统筹/路演）", "研发：__________（前端/后端/数据）", "运营：__________（内容/活动/渠道）", "设计：__________（视觉/交互/素材）"])
    add_section_card(slide, Inches(6.95), Inches(1.35), Inches(5.9), Inches(5.5), "优势与资源（填空）", ["对接资源：茶园/非遗工坊/景区/供应链/渠道", "技术能力：Web开发、数据分析、AI集成", "商业能力：运营增长、品牌策划、团购拓展"])
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    bg2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = col_primary
    bg2.line.fill.background()
    if img_mascot.exists():
        slide.shapes.add_picture(str(img_mascot), prs.slide_width - Inches(4.0), Inches(2.2), width=Inches(3.1))
    box2 = slide.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(9.8), Inches(2.8))
    tf2 = box2.text_frame
    tf2.clear()
    p = tf2.paragraphs[0]
    p.text = "谢谢评委老师"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p2 = tf2.add_paragraph()
    p2.text = "Q & A"
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(30)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(0xE6, 0xFF, 0xEF)
    p2.space_before = Pt(10)
    slides.append(slide)

    total = len(slides)
    for i, s in enumerate(slides, start=1):
        if i in (1, total):
            continue
        foot = s.shapes.add_textbox(prs.slide_width - Inches(1.2), prs.slide_height - Inches(0.5), Inches(1.0), Inches(0.3))
        fp = foot.text_frame.paragraphs[0]
        fp.text = f"{i}/{total}"
        fp.font.name = "Microsoft YaHei"
        fp.font.size = Pt(10)
        fp.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        fp.alignment = PP_ALIGN.RIGHT

    prs.save(str(out))
    print(out)


if __name__ == "__main__":
    main()
